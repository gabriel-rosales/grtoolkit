import sys, os
from pptx import Presentation
from grtoolkit.Storage import File
from grtoolkit.PowerPoint import duplicate_slide

parent = os.getcwd()
QUOTES = f'{parent}\\Quotes.txt'
OLD_PPT = f'{parent}\\QUOTE_PPT_TEMPLATE.pptm'
NEW_PPT = f'{parent}\\QUOTE_NEW.pptm'
prs = Presentation(OLD_PPT)

QUOTESContent = File(QUOTES, enc=1).read()
QUOTESList = QUOTESContent.split("\n\n")

#DUPLICATE SLIDE 0 FROM TEMPLATE
duplicate_slide(prs,0,len(QUOTESList)-1)
#CHANGE TEXT IN ALL SLIDES AS PER NOTEPAD FILE
i=0
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            shape.text_frame.clear()
            p=shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = QUOTESList[i]
            font=run.font
            font.name='Gabriola'
    i=i+1
prs.save(NEW_PPT)